import os
import sys
import zipfile
from pptx import Presentation
import xml.etree.ElementTree as ET

# Ensure UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

SRC_PATH = r"C:\Users\shri hari computer\Downloads\Meerut Media Plan_Master Data.pptx"
OUT_DIR = r"C:\Users\shri hari computer\Downloads"

def split_into_8_parts():
    print(f"[*] Loading original presentation: {SRC_PATH}")
    prs_master = Presentation(SRC_PATH)
    total_slides = len(prs_master.slides)
    print(f"[*] Total Slides found: {total_slides}")

    # Determine 8 balanced chunks
    num_parts = 8
    chunk_size = total_slides // num_parts
    remainder = total_slides % num_parts

    ranges = []
    current_start = 0
    for p in range(num_parts):
        # Distribute remainder across first few parts
        extra = 1 if p < remainder else 0
        current_end = current_start + chunk_size + extra
        ranges.append((current_start, current_end))
        current_start = current_end

    generated_files = []

    for part_idx, (start_idx, end_idx) in enumerate(ranges, 1):
        num_in_part = end_idx - start_idx
        slide_range_str = f"Slides {start_idx + 1} to {end_idx} ({num_in_part} slides)"
        out_filename = f"Meerut_Media_Plan_Part_{part_idx}_of_8.pptx"
        out_filepath = os.path.join(OUT_DIR, out_filename)
        
        print(f"\n[+] Generating Part {part_idx}/8: {slide_range_str}...")
        
        prs = Presentation(SRC_PATH)
        total = len(prs.slides)

        # Remove slides after end_idx (from back to front)
        for i in range(total - 1, end_idx - 1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # Remove slides before start_idx (from start_idx-1 down to 0)
        for i in range(start_idx - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        temp_out = os.path.join(OUT_DIR, f"temp_part_{part_idx}.pptx")
        prs.save(temp_out)

        # Optimize by removing unused media from the zip package
        clean_zip_media(temp_out, out_filepath)
        if os.path.exists(temp_out):
            os.remove(temp_out)

        file_size_mb = os.path.getsize(out_filepath) / (1024 * 1024)
        print(f"[OK] Created: {out_filename} | Size: {file_size_mb:.2f} MB | {num_in_part} Slides")
        generated_files.append({
            'part': part_idx,
            'name': out_filename,
            'path': out_filepath,
            'slides': num_in_part,
            'range': f"{start_idx + 1}-{end_idx}",
            'size_mb': file_size_mb
        })

    # Remove temporary test file if exists
    test_file = os.path.join(OUT_DIR, "test_part1.pptx")
    if os.path.exists(test_file):
        os.remove(test_file)

    print("\n" + "="*70)
    print("ALL 8 PARTS CREATED SUCCESSFULLY IN DOWNLOADS FOLDER:")
    print("="*70)
    for item in generated_files:
        print(f"Part {item['part']}: {item['name']} | Range: Slides {item['range']} | {item['slides']} Slides | {item['size_mb']:.2f} MB")
    print("="*70)

def clean_zip_media(src_pptx, dst_pptx):
    """
    Scans slide relationships in PPTX and only keeps referenced media files in the archive.
    """
    with zipfile.ZipFile(src_pptx, 'r') as zin:
        referenced_media = set()
        for filename in zin.namelist():
            if filename.startswith('ppt/slides/_rels/') and filename.endswith('.rels'):
                content = zin.read(filename).decode('utf-8', errors='ignore')
                try:
                    root = ET.fromstring(content)
                    for rel in root:
                        target = rel.attrib.get('Target', '')
                        if 'media/' in target:
                            media_name = os.path.basename(target)
                            referenced_media.add(f"ppt/media/{media_name}")
                except Exception:
                    pass

        with zipfile.ZipFile(dst_pptx, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename.startswith('ppt/media/'):
                    if item.filename in referenced_media:
                        zout.writestr(item, zin.read(item.filename))
                else:
                    zout.writestr(item, zin.read(item.filename))

if __name__ == '__main__':
    split_into_8_parts()
